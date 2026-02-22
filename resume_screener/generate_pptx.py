"""Generate a PowerPoint (.pptx) from the Resume Screener AI presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── Colour constants ──
BG       = RGBColor(0x0A, 0x0A, 0x0F)
SURFACE  = RGBColor(0x16, 0x19, 0x22)
TEXT_PRI = RGBColor(0xF0, 0xF0, 0xF5)
TEXT_SEC = RGBColor(0x8B, 0x8F, 0xA3)
TEXT_MUT = RGBColor(0x5A, 0x5E, 0x72)
ACCENT   = RGBColor(0x63, 0x66, 0xF1)
GREEN    = RGBColor(0x22, 0xC5, 0x5E)
RED      = RGBColor(0xEF, 0x44, 0x44)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
CYAN     = RGBColor(0x06, 0xB6, 0xD4)
PINK     = RGBColor(0xEC, 0x48, 0x99)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # Blank

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height, size=18, bold=False,
             color=TEXT_PRI, align=PP_ALIGN.CENTER, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_badge(slide, text, top=0.8, color=ACCENT):
    add_text(slide, f"● {text}", 4.5, top, 4.3, 0.4, size=11, bold=True, color=color)

def add_card_row(slide, cards, top, cols=3):
    """cards = list of (emoji, title, desc)"""
    gap = 0.3
    total_w = 12.0
    card_w = (total_w - gap * (cols - 1)) / cols
    left_start = (13.333 - total_w) / 2
    for i, (emoji, title, desc) in enumerate(cards):
        col = i % cols
        row = i // cols
        x = left_start + col * (card_w + gap)
        y = top + row * 1.6
        # card background
        shape = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(card_w), Inches(1.4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.fill.background()
        # text
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{emoji} {title}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_PRI
        p.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_SEC
        p2.font.name = "Calibri"

def add_code_block(slide, code, top, height=2.5):
    shape = slide.shapes.add_shape(
        1, Inches(1.5), Inches(top), Inches(10.3), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_SEC
    p.font.name = "Courier New"
    p.alignment = PP_ALIGN.LEFT

# ═══════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_text(slide, "Resume Screener AI", 1.5, 2.0, 10.3, 1.2, size=44, bold=True, color=ACCENT)
add_text(slide, "Intelligent Resume-to-Job-Description matching\npowered by Google Gemini AI", 2.5, 3.4, 8.3, 1.0, size=18, color=TEXT_SEC)
add_text(slide, "━━━━━━━━━", 5.5, 4.5, 2.3, 0.3, size=14, color=CYAN)
add_text(slide, "QUANTUM BITS", 4.5, 5.0, 4.3, 0.5, size=14, bold=True, color=TEXT_MUT)

# ═══════════════════════════════════════════════════
# SLIDE 2: Problem Statement
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "PROBLEM", color=RED)
add_text(slide, "The Hiring Challenge", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
# Stats
for i, (val, lbl, clr) in enumerate([
    ("250+", "Resumes per job posting", RED),
    ("23 sec", "Avg. time per manual review", AMBER),
    ("75%", "Resumes are unqualified", PINK),
]):
    x = 1.5 + i * 3.8
    add_text(slide, val, x, 2.8, 3.3, 0.8, size=40, bold=True, color=clr)
    add_text(slide, lbl, x, 3.7, 3.3, 0.5, size=12, color=TEXT_MUT)
add_text(slide, "Manual screening is slow, biased, and inconsistent — costing companies time and top talent.",
         2, 4.8, 9.3, 0.8, size=16, color=TEXT_SEC)

# ═══════════════════════════════════════════════════
# SLIDE 3: Solution
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "SOLUTION", color=GREEN)
add_text(slide, "AI-Powered Resume Screening", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "Upload a resume and job description → Get an instant match score,\nskill gap analysis, and a summary — all in seconds.",
         2, 2.4, 9.3, 0.8, size=16, color=TEXT_SEC)
add_card_row(slide, [
    ("📄", "Upload PDF", "Resume & JD in PDF or text format"),
    ("🤖", "AI Analysis", "Gemini LLM evaluates the match"),
    ("📊", "Score & Report", "Match %, matching & missing skills"),
], top=3.8)

# ═══════════════════════════════════════════════════
# SLIDE 4: Tech Stack
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "TECHNOLOGY", color=CYAN)
add_text(slide, "Tech Stack", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "A modern full-stack architecture combining Python and JavaScript",
         2.5, 2.4, 8.3, 0.6, size=16, color=TEXT_SEC)
techs = ["🐍 Python", "⚡ FastAPI", "🦜 LangChain", "✨ Gemini AI",
         "⚛️ React", "🚀 Vite", "📜 PyPDF2", "🎨 CSS"]
add_card_row(slide, [(t.split(" ")[0], " ".join(t.split(" ")[1:]), "") for t in techs], top=3.5, cols=4)

# ═══════════════════════════════════════════════════
# SLIDE 5: Architecture
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "ARCHITECTURE", color=ACCENT)
add_text(slide, "System Architecture", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
flow = "👤 User  →  ⚛️ React Frontend  →  🔗 REST API  →  ⚡ FastAPI  →  📄 PDF Parser  →  🤖 Gemini LLM  →  📊 JSON Result"
add_text(slide, flow, 1, 3.0, 11.3, 1.0, size=16, bold=True, color=TEXT_SEC)
add_text(slide, "End-to-end pipeline from resume upload to AI-powered scoring",
         2.5, 4.5, 8.3, 0.6, size=16, color=TEXT_MUT)

# ═══════════════════════════════════════════════════
# SLIDE 6: Backend — FastAPI
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "BACKEND", color=GREEN)
add_text(slide, "FastAPI Server", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "High-performance Python framework with automatic API documentation",
         2.5, 2.3, 8.3, 0.6, size=16, color=TEXT_SEC)
code = """from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware

app = FastAPI(title="Resume Screening AI")

@app.post("/analyze")
async def analyze(
    resume: UploadFile = File(...),
    job_description: str = Form(None),
    job_description_file: UploadFile = File(None)
):
    # Parse → Analyze → Return JSON"""
add_code_block(slide, code, top=3.2, height=3.0)

# ═══════════════════════════════════════════════════
# SLIDE 7: PDF Parser
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "MODULE", color=AMBER)
add_text(slide, "PDF Parsing Engine", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "Extracts raw text from uploaded PDF files using PyPDF2",
         2.5, 2.3, 8.3, 0.6, size=16, color=TEXT_SEC)
code = """import PyPDF2, io

def extract_text_from_pdf(file_content: bytes) -> str:
    \"\"\"Extracts text from a PDF file content.\"\"\"
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() or ""
    return text.strip()"""
add_code_block(slide, code, top=3.2, height=2.5)
add_card_row(slide, [
    ("📥", "Input", "Raw PDF bytes from file upload"),
    ("📤", "Output", "Clean plain text for LLM analysis"),
], top=6.0, cols=2)

# ═══════════════════════════════════════════════════
# SLIDE 8: LLM Logic
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "AI CORE", color=PINK)
add_text(slide, "LLM Scoring Engine", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "LangChain orchestrates the prompt → Gemini → structured output pipeline",
         2, 2.3, 9.3, 0.6, size=16, color=TEXT_SEC)
code = """class ResumeScore(BaseModel):
    score: int             # Match 0–100
    matching_skills: List[str]  # Skills found in both
    missing_skills: List[str]   # Required but absent
    summary: str            # Brief suitability note

llm = ChatGoogleGenerativeAI(model="gemini-flash-latest")
parser = JsonOutputParser(pydantic_object=ResumeScore)

chain = prompt | llm | parser  # LangChain pipeline"""
add_code_block(slide, code, top=3.2, height=2.8)

# ═══════════════════════════════════════════════════
# SLIDE 9: Pydantic Schema
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "DATA MODEL", color=CYAN)
add_text(slide, "Structured Output", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "Pydantic model ensures type-safe, validated JSON output from the LLM",
         2, 2.3, 9.3, 0.6, size=16, color=TEXT_SEC)
add_card_row(slide, [
    ("📊", "Score", "0–100 integer match percentage"),
    ("✅", "Matching Skills", "Skills present in both resume & JD"),
    ("❌", "Missing Skills", "Required skills not found in resume"),
    ("📝", "Summary", "Brief AI analysis of fit"),
], top=3.2, cols=4)
code = """// Example API Response
{
  "score": 82,
  "matching_skills": ["Python", "React", "SQL"],
  "missing_skills": ["Docker", "AWS"],
  "summary": "Strong match with minor gaps..."
}"""
add_code_block(slide, code, top=5.0, height=2.0)

# ═══════════════════════════════════════════════════
# SLIDE 10: Frontend Overview
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "FRONTEND", color=ACCENT)
add_text(slide, "React + Vite UI", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "Fast, responsive single-page application with modern design",
         2.5, 2.3, 8.3, 0.6, size=16, color=TEXT_SEC)
add_card_row(slide, [
    ("⚛️", "React 19", "Latest React with hooks for state management"),
    ("🚀", "Vite 7", "Lightning-fast HMR dev server and optimized bundling"),
    ("🎨", "Glassmorphism CSS", "Custom design system with CSS variables"),
    ("📡", "Axios", "Clean HTTP client for FormData uploads"),
], top=3.3, cols=2)

# ═══════════════════════════════════════════════════
# SLIDE 11: UI Components
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "UI", color=GREEN)
add_text(slide, "User Interface Components", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_card_row(slide, [
    ("📤", "Resume Upload", "Drag & drop zone for PDF resumes with file validation"),
    ("📝", "JD Input", "Toggle between pasting text and uploading a PDF"),
    ("🔘", "Analyze CTA", "Full-width button with loading spinner animation"),
    ("🎯", "Score Ring", "Animated SVG circle showing match percentage"),
    ("📊", "Skills Tags", "Green/red tagged chips for matching & missing skills"),
    ("📋", "Summary", "AI-generated paragraph about candidate suitability"),
], top=2.5, cols=3)

# ═══════════════════════════════════════════════════
# SLIDE 12: Score Visualization
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "FEATURE", color=GREEN)
add_text(slide, "Score Visualization", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "Dynamic SVG ring with color-coded match scores",
         2.5, 2.3, 8.3, 0.6, size=16, color=TEXT_SEC)
add_text(slide, "82%", 5.2, 3.0, 3.0, 1.5, size=64, bold=True, color=GREEN)
add_text(slide, "MATCH", 5.5, 4.3, 2.3, 0.5, size=14, color=TEXT_MUT)
add_card_row(slide, [
    ("🟢", "≥ 75%", "Green — Strong match"),
    ("🟡", "50–74%", "Amber — Moderate fit"),
    ("🔴", "< 50%", "Red — Weak match"),
], top=5.5, cols=3)

# ═══════════════════════════════════════════════════
# SLIDE 13: API Design
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "API", color=AMBER)
add_text(slide, "RESTful API Design", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "Single endpoint powered by FastAPI with automatic OpenAPI documentation",
         2, 2.3, 9.3, 0.6, size=16, color=TEXT_SEC)
code = """# Endpoint
POST /analyze

# Request (multipart/form-data)
  resume:               File   (required, PDF)
  job_description:      string (optional, text)
  job_description_file: File   (optional, PDF)

# Response (application/json)
  { score, matching_skills, missing_skills, summary }"""
add_code_block(slide, code, top=3.2, height=2.5)
add_card_row(slide, [
    ("🔀", "CORS Enabled", "Allows cross-origin requests from any frontend"),
    ("📖", "Auto Docs", "Swagger UI at /docs — test the API interactively"),
], top=6.0, cols=2)

# ═══════════════════════════════════════════════════
# SLIDE 14: LangChain Pipeline
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "PIPELINE", color=PINK)
add_text(slide, "LangChain Pipeline", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "A composable chain: Prompt → LLM → Parser",
         2.5, 2.3, 8.3, 0.6, size=16, color=TEXT_SEC)
add_card_row(slide, [
    ("📝", "Prompt Template", "Injects resume + JD + format instructions"),
    ("🤖", "Gemini Flash", "Processes & evaluates match"),
    ("📊", "JSON Parser", "Validates into Pydantic model"),
], top=3.2, cols=3)
code = """chain = prompt | llm | parser
result = chain.invoke({
    "resume_text": resume_text,
    "job_description": job_description,
    "format_instructions": parser.get_format_instructions()
})"""
add_code_block(slide, code, top=5.2, height=1.8)

# ═══════════════════════════════════════════════════
# SLIDE 15: Dual Input Modes
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "FLEXIBILITY", color=CYAN)
add_text(slide, "Dual JD Input Modes", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "Users can provide job descriptions in two ways",
         2.5, 2.3, 8.3, 0.6, size=16, color=TEXT_SEC)
add_card_row(slide, [
    ("📝", "Paste Text", "Quick for short JDs\nCopy from job boards\nInstant — no file needed\nEdit before analysis"),
    ("📄", "Upload PDF", "Full formatted JDs\nPreserves document layout\nMulti-page support\nSame PyPDF2 engine"),
], top=3.3, cols=2)

# ═══════════════════════════════════════════════════
# SLIDE 16: Error Handling
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "ROBUSTNESS", color=RED)
add_text(slide, "Error Handling", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "Comprehensive validation at every layer",
         2.5, 2.3, 8.3, 0.6, size=16, color=TEXT_SEC)
add_card_row(slide, [
    ("🔒", "File Validation", "Only PDF files accepted — non-PDF uploads rejected with 400"),
    ("📋", "Missing Input", "JD is required — clear error if neither text nor file provided"),
    ("🖥️", "Frontend Guards", "Client-side validation prevents premature submission"),
    ("🔎", "Server Traceback", "Full stack traces logged on failure for fast debugging"),
], top=3.3, cols=2)

# ═══════════════════════════════════════════════════
# SLIDE 17: Design System
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "DESIGN", color=ACCENT)
add_text(slide, "Design System", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "Clean, minimal UI with CSS variables and responsive layout",
         2.5, 2.3, 8.3, 0.6, size=16, color=TEXT_SEC)
add_card_row(slide, [
    ("🎨", "Slate Palette", "10-shade slate color scale"),
    ("✨", "Micro Animations", "slideUp, spinner, ring transitions"),
    ("📱", "Responsive", "Breakpoint at 600px — stacks layouts"),
    ("🔤", "Inter Font", "Google Fonts with weights 400–700"),
    ("📐", "Layout", "Flexbox with 640px max-width container"),
    ("💫", "Hover Effects", "Border color transitions, button shifts"),
], top=3.3, cols=3)

# ═══════════════════════════════════════════════════
# SLIDE 18: Deployment
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "DEPLOYMENT", color=GREEN)
add_text(slide, "Hosting Strategy", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_text(slide, "Flexible deployment options for production",
         2.5, 2.3, 8.3, 0.6, size=16, color=TEXT_SEC)
add_card_row(slide, [
    ("🖥️", "Backend Hosting", "Render / Railway — FastAPI with uvicorn"),
    ("🌐", "Frontend Hosting", "Vercel / Netlify — Vite dist folder"),
    ("🔗", "Unified Mode", "Backend serves frontend dist — single server"),
    ("🚇", "Dev Tunneling", "ngrok for instant public URLs during dev"),
], top=3.3, cols=2)

# ═══════════════════════════════════════════════════
# SLIDE 19: Key Features
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "SUMMARY", color=ACCENT)
add_text(slide, "Key Features", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
features = [
    "✓  Instant AI-powered resume-to-JD matching with 0–100 scoring",
    "✓  Skill gap analysis — see exactly which skills match and which are missing",
    "✓  Dual input modes — paste text or upload PDF for job descriptions",
    "✓  Beautiful responsive UI with animations and glassmorphism design",
    "✓  Structured JSON output validated by Pydantic models",
    "✓  LangChain pipeline with Google Gemini Flash for fast inference",
    "✓  Production-ready with CORS, error handling, and env config",
]
for i, feat in enumerate(features):
    add_text(slide, feat, 2, 2.5 + i * 0.6, 9.3, 0.5, size=16,
             color=TEXT_SEC, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════
# SLIDE 20: Future Scope
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_badge(slide, "ROADMAP", color=AMBER)
add_text(slide, "Future Scope", 1.5, 1.4, 10.3, 0.8, size=36, bold=True)
add_card_row(slide, [
    ("📦", "Batch Processing", "Upload multiple resumes and rank them against a single JD"),
    ("🗃️", "Database Storage", "Store past analyses for comparison and tracking"),
    ("🔐", "Authentication", "User accounts for recruiters to manage pipelines"),
    ("📈", "Analytics Dashboard", "Visualize hiring trends, skill demands, and pools"),
    ("💡", "Resume Improvement", "AI suggestions for candidates to improve scores"),
    ("🔌", "ATS Integration", "Connect with existing Applicant Tracking Systems"),
], top=2.5, cols=3)

# ═══════════════════════════════════════════════════
# SLIDE 21: Thank You
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)
add_text(slide, "Thank You", 1.5, 2.2, 10.3, 1.2, size=48, bold=True, color=ACCENT)
add_text(slide, "━━━━━━━━━", 5.5, 3.5, 2.3, 0.3, size=14, color=CYAN)
add_text(slide, "Resume Screener AI — by Quantum Bits", 2.5, 4.0, 8.3, 0.6, size=18, color=TEXT_SEC)
add_text(slide, "Built with FastAPI • React • LangChain • Google Gemini",
         2.5, 4.8, 8.3, 0.6, size=14, color=TEXT_MUT)

# ── Save ──
out_path = os.path.join(os.path.dirname(__file__), "Resume_Screener_AI_Presentation.pptx")
prs.save(out_path)
print(f"✅ Saved to: {out_path}")
